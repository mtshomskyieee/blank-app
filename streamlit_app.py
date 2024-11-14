style_css = """
/* Main container styling */
.stApp {
    max-width: 1200px;
    margin: 0 auto;
}

/* Header styling */
.stTitle {
    color: #1E1E1E;
    text-align: center;
    padding: 1rem 0;
}

/* Form styling */
.stForm {
    background-color: #f8f9fa;
    padding: 2rem;
    border-radius: 10px;
    margin-bottom: 2rem;
}

/* Button styling */
.stButton > button {
    width: 100%;
    border-radius: 5px;
    padding: 0.5rem 1rem;
}

/* Slide container styling */
.slide-container {
    background-color: white;
    padding: 2rem;
    border-radius: 10px;
    box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    margin: 2rem 0;
}

/* Navigation buttons */
.stButton > button:first-child {
    background-color: #007bff;
    color: white;
}

/* Progress bar styling */
.stProgress > div {
    height: 5px;
    background-color: #007bff;
}

/* Image styling */
.stImage {
    margin: 1rem 0;
    border-radius: 5px;
}

/* Error message styling */
.stAlert {
    padding: 1rem;
    border-radius: 5px;
    margin: 1rem 0;
}

"""
# Save string to file
#with open('styles.css', 'w') as file:
#    file.write(styles_css)


import requests
from PIL import Image
from io import BytesIO

def process_image_response(image_url):
    try:
        response = requests.get(image_url)
        if response.status_code != 200:
            print(f"Failed to fetch image: HTTP {response.status_code}")
            return None
            
        image = Image.open(BytesIO(response.content))
        
        # Convert to RGB if image is in RGBA
        if image.mode == 'RGBA':
            image = image.convert('RGB')
            
        # Create BytesIO object to store processed image
        img_byte_arr = BytesIO()
        image.save(img_byte_arr, format='JPEG')
        img_byte_arr = img_byte_arr.getvalue()
        
        return img_byte_arr
    except requests.RequestException as e:
        print(f"Network error while fetching image: {e}")
        return None
    except Image.UnidentifiedImageError as e:
        print(f"Invalid image data received: {e}")
        return None
    except Exception as e:
        print(f"Error processing image: {e}")
        return None
    

class SlideTemplates:
    @staticmethod
    def get_template(style):
        templates = {
            "Professional": {
                "title_format": "### {title}",
                "content_format": "{content}",
                "image_position": "top"
            },
            "Creative": {
                "title_format": "## 🎨 {title}",
                "content_format": "> {content}",
                "image_position": "bottom"
            },
            "Academic": {
                "title_format": "### 📚 {title}",
                "content_format": "{content}",
                "image_position": "right"
            },
            "Casual": {
                "title_format": "## ✨ {title}",
                "content_format": "{content}",
                "image_position": "left"
            }
        }
        return templates.get(style, templates["Professional"])
import os
from langchain_openai import ChatOpenAI
from langchain.prompts import ChatPromptTemplate
import openai

import json

class PresentationGenerator:
    def __init__(self):
        self.llm = ChatOpenAI(
            temperature=0.7,
            model_name="gpt-4",  # Using GPT-4 for better content generation
            request_timeout=60
        )
        self.client = openai.OpenAI()  # Initialize OpenAI client
        
    def generate_slide_content(self, topic, slide_number, total_slides, style, slide_outline=None):
        prompt = ChatPromptTemplate.from_template('''
Create a presentation slide with the following requirements:
- Slide number: {slide_number} of {total_slides}
- Topic: {topic}
- Style: {style}
- Slide Outline: {slide_outline}

Follow the provided outline point for this slide number. Make sure the content aligns with the outline point.

Return the response in strict JSON format without any additional text or formatting.
Use this exact structure:
{{"title": "Slide Title", "content": "Main Content", "image_prompt": "Image Description"}}
''')
        
        chain = prompt | self.llm
        response = chain.invoke({
            "topic": topic,
            "slide_number": slide_number,
            "total_slides": total_slides,
            "style": style,
            "slide_outline": slide_outline if slide_outline else "Generate appropriate content for this slide position"
        })
        
        try:
            # Extract content from ChatMessage
            response_str = response.content if hasattr(response, 'content') else str(response)
            
            # Clean the response string
            response_str = response_str.strip()
            response_str = ' '.join(response_str.split())
            response_str = response_str.replace('```json', '').replace('```', '')
            
            # Parse JSON
            slide_content = json.loads(response_str)
            
            # Validate required keys
            required_keys = ['title', 'content', 'image_prompt']
            missing_keys = [key for key in required_keys if key not in slide_content]
            if missing_keys:
                raise ValueError(f"Missing required keys: {missing_keys}")
            
            return slide_content
        except json.JSONDecodeError as e:
            raise ValueError(f"Failed to parse LLM response as JSON. Response: {response_str}. Error: {str(e)}")
        except Exception as e:
            raise ValueError(f"Error processing slide content: {str(e)}")
    
    def generate_outline(self, topic, num_slides):
        prompt = ChatPromptTemplate.from_template('''
Create a presentation outline for the following topic:
- Topic: {topic}
- Number of slides: {num_slides}

Generate a concise outline with exactly {num_slides} points that would make an effective presentation.
Each point should be a brief description of what that slide should cover.

Return the response as a JSON array of strings, where each string is a slide outline point.
Example format: ["Introduction to topic", "Key point 1", "Key point 2", "Conclusion"]
''')
        
        chain = prompt | self.llm
        response = chain.invoke({
            "topic": topic,
            "num_slides": num_slides
        })
        
        try:
            response_str = response.content if hasattr(response, 'content') else str(response)
            response_str = response_str.strip()
            response_str = ' '.join(response_str.split())
            response_str = response_str.replace('```json', '').replace('```', '')
            
            outline = json.loads(response_str)
            if not isinstance(outline, list) or len(outline) != num_slides:
                raise ValueError("Invalid outline format")
            
            return outline
        except Exception as e:
            raise ValueError(f"Error generating outline: {str(e)}")
    
    def generate_image(self, prompt):
        try:
            response = self.client.images.generate(
                model="dall-e-2",
                prompt=prompt,
                n=1,
                size="512x512"
            )
            return process_image_response(response.data[0].url)
        except Exception as e:
            print(f"Error generating image: {e}")
            return None
    
    def generate_presentation(self, topic, num_slides, style, outline=None):
        if outline is None:
            outline = self.generate_outline(topic, num_slides)
            
        slides = []
        
        for i in range(num_slides):
            try:
                print(f"Generating slide {i+1}/{num_slides}...")  # Progress update
                # Generate slide content using the outline point
                slide_content = self.generate_slide_content(
                    topic, 
                    i+1, 
                    num_slides, 
                    style, 
                    outline[i]
                )
                
                print(f"Generating image for slide {i+1}...")  # Progress update
                # Generate image for the slide
                image = self.generate_image(slide_content['image_prompt'])
                
                # Combine content and image
                slide = {
                    'title': slide_content['title'],
                    'content': slide_content['content'],
                    'image': image
                }
                
                slides.append(slide)
            except Exception as e:
                raise ValueError(f"Error generating slide {i+1}: {str(e)}")
        
        return slides

    def export_to_pptx(self, slides, filename="presentation.pptx"):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from io import BytesIO
            
            # Create presentation
            prs = Presentation()
            
            # Define consistent dimensions
            left_margin = Inches(0.5)
            top_margin = Inches(1.5)
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Calculate content area dimensions
            content_width = slide_width - (2 * Inches(0.5))  # Subtract margins
            content_height = slide_height - Inches(1.5) - Inches(0.5)  # Subtract top and bottom margins
            
            # Calculate column widths (divide space equally)
            column_width = content_width / 2
            
            for slide in slides:
                # Add a slide with blank layout
                layout = prs.slide_layouts[6]  # Blank layout
                slide_obj = prs.slides.add_slide(layout)
                
                # Add title
                title_box = slide_obj.shapes.add_textbox(
                    left_margin, 
                    Inches(0.5), 
                    content_width, 
                    Inches(1)
                )
                # Add title with line breaks
                title_frame = title_box.text_frame
                title_text = slide['title']

                # Add line breaks every 42 characters
                if len(title_text) > 42:
                    # Split into chunks of 42 characters at word boundaries
                    words = title_text.split()
                    current_line = []
                    current_length = 0
                    
                    formatted_lines = []
                    for word in words:
                        if current_length + len(word) + 1 <= 42:  # +1 for space
                            current_line.append(word)
                            current_length += len(word) + 1
                        else:
                            formatted_lines.append(' '.join(current_line))
                            current_line = [word]
                            current_length = len(word)
                    
                    if current_line:
                        formatted_lines.append(' '.join(current_line))
                    
                    # Join lines with newline character
                    title_frame.text = '\n'.join(formatted_lines)
                else:
                    title_frame.text = title_text

                title_frame.paragraphs[0].font.size = Pt(32)
                title_frame.paragraphs[0].font.bold = True
                
                # Add text content on the left side
                text_box = slide_obj.shapes.add_textbox(
                    left_margin,
                    top_margin,
                    column_width,
                    content_height
                )
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_frame.text = slide['content']
                
                # Add image on the right side if available
                if slide.get('image'):
                    img_stream = BytesIO(slide['image'])
                    # Position image on the right side
                    image_left = left_margin + column_width + Inches(0.25)  # Add small gap
                    image_width = column_width - Inches(0.5)  # Slightly smaller than column
                    image_height = content_height - Inches(0.5)
                    
                    slide_obj.shapes.add_picture(
                        img_stream,
                        image_left,
                        top_margin,
                        width=image_width,
                        height=image_height
                    )
            
            # Save presentation
            prs.save(filename)
            return filename
        except Exception as e:
            raise ValueError(f"Error exporting to PowerPoint: {str(e)}")
import streamlit as st

import os

# Page configuration
st.set_page_config(
    page_title="AI Presentation Generator",
    page_icon="🎯",
    layout="wide"
)

# Custom CSS
st.markdown("""
<style>
.stApp {
    max-width: 1200px;
    margin: 0 auto;
}
.stTitle {
    color: #1E1E1E;
    text-align: center;
    padding: 1rem 0;
}
.stForm {
    background-color: #f8f9fa;
    padding: 2rem;
    border-radius: 10px;
    margin-bottom: 2rem;
}
.stButton > button {
    width: 100%;
    border-radius: 5px;
    padding: 0.5rem 1rem;
}
.slide-container {
    background-color: white;
    padding: 2rem;
    border-radius: 10px;
    box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    margin: 2rem 0;
}
.stProgress > div {
    height: 5px;
    background-color: #007bff;
}
.stImage {
    margin: 1rem 0;
    border-radius: 5px;
}
.stAlert {
    padding: 1rem;
    border-radius: 5px;
    margin: 1rem 0;
}
</style>
""", unsafe_allow_html=True)

def main():
    st.title("🎯 AI Presentation Generator")
    
    # Initialize session state
    if 'current_slide' not in st.session_state:
        st.session_state.current_slide = 0
    if 'slides' not in st.session_state:
        st.session_state.slides = []
    if 'outline' not in st.session_state:
        st.session_state.outline = None
    
    # Input section
    with st.form("presentation_form"):
        topic = st.text_input("Enter your presentation topic:", 
                            placeholder="e.g., The Future of Artificial Intelligence")
        num_slides = st.slider("Number of slides:", min_value=3, max_value=10, value=5)
        style = st.selectbox("Presentation style:", 
                           ["Professional", "Creative", "Academic", "Casual"])
        
        # Generate outline button
        generate_outline = st.form_submit_button("Generate Outline")
        
        if generate_outline and topic:
            with st.spinner("Generating presentation outline..."):
                try:
                    generator = PresentationGenerator()
                    st.session_state.outline = generator.generate_outline(topic, num_slides)
                except Exception as e:
                    st.error(f"An error occurred: {str(e)}")
    
    # Display and edit outline
    if st.session_state.outline:
        st.subheader("Presentation Outline")
        st.info("Edit the outline points below to customize your presentation structure:")
        
        edited_outline = []
        for i, point in enumerate(st.session_state.outline):
            edited_point = st.text_input(f"Slide {i+1}", value=point, key=f"outline_{i}")
            edited_outline.append(edited_point)
        
        # Generate presentation button
        if st.button("Generate Presentation with Custom Outline"):
            with st.spinner("Generating your presentation..."):
                try:
                    generator = PresentationGenerator()
                    st.session_state.slides = generator.generate_presentation(
                        topic, num_slides, style, edited_outline
                    )
                    st.session_state.current_slide = 0
                    st.success("Presentation generated successfully!")
                except Exception as e:
                    st.error(f"An error occurred: {str(e)}")
    
    # Display presentation
    if st.session_state.slides:
        col1, col2, col3 = st.columns([1, 3, 1])
        
        with col1:
            if st.button("⬅️ Previous") and st.session_state.current_slide > 0:
                st.session_state.current_slide -= 1
                
        with col3:
            if st.button("Next ➡️") and st.session_state.current_slide < len(st.session_state.slides) - 1:
                st.session_state.current_slide += 1
        
        # Display current slide
        current_slide = st.session_state.slides[st.session_state.current_slide]
        
        with st.container():
            st.markdown(f"### Slide {st.session_state.current_slide + 1}/{len(st.session_state.slides)}")
            
            # Display slide content
            st.markdown(f"## {current_slide['title']}")
            if 'image' in current_slide and current_slide['image'] is not None:
                st.image(current_slide['image'], use_container_width=True)
            else:
                st.warning("Image generation failed. Displaying text content only.")
            st.markdown(current_slide['content'])
            
            # Progress bar
            progress = (st.session_state.current_slide + 1) / len(st.session_state.slides)
            st.progress(progress)

            # Add export button
            if st.button("Export to PowerPoint"):
                try:
                    with st.spinner("Exporting presentation..."):
                        generator = PresentationGenerator()
                        filename = generator.export_to_pptx(st.session_state.slides)
                        
                        # Read file for download
                        with open(filename, "rb") as file:
                            btn = st.download_button(
                                label="Download Presentation",
                                data=file,
                                file_name="presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        
                        # Clean up file after download button is created
                        if os.path.exists(filename):
                            os.remove(filename)
                            
                except Exception as e:
                    st.error(f"Error exporting presentation: {str(e)}")

if __name__ == "__main__":
    main()
